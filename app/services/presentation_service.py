import io
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

THEMES = {
    "teal": {
        "primary": RGBColor(0, 137, 123),
        "secondary": RGBColor(0, 77, 64),
        "accent": RGBColor(38, 166, 154),
        "bg": RGBColor(240, 253, 244),
        "text_dark": RGBColor(30, 41, 59),
        "text_light": RGBColor(255, 255, 255)
    },
    "blue": {
        "primary": RGBColor(30, 136, 229),
        "secondary": RGBColor(13, 71, 161),
        "accent": RGBColor(66, 165, 245),
        "bg": RGBColor(239, 246, 255),
        "text_dark": RGBColor(30, 41, 59),
        "text_light": RGBColor(255, 255, 255)
    },
    "purple": {
        "primary": RGBColor(142, 36, 170),
        "secondary": RGBColor(74, 20, 140),
        "accent": RGBColor(171, 71, 188),
        "bg": RGBColor(250, 245, 255),
        "text_dark": RGBColor(30, 41, 59),
        "text_light": RGBColor(255, 255, 255)
    },
    "amber": {
        "primary": RGBColor(217, 119, 6),
        "secondary": RGBColor(120, 53, 15),
        "accent": RGBColor(245, 158, 11),
        "bg": RGBColor(255, 251, 235),
        "text_dark": RGBColor(30, 41, 59),
        "text_light": RGBColor(255, 255, 255)
    },
    "rose": {
        "primary": RGBColor(225, 29, 72),
        "secondary": RGBColor(136, 19, 55),
        "accent": RGBColor(244, 63, 94),
        "bg": RGBColor(255, 241, 242),
        "text_dark": RGBColor(30, 41, 59),
        "text_light": RGBColor(255, 255, 255)
    }
}

def generate_pptx_presentation(
    client_name: str,
    period: str,
    theme: str = "teal",
    num_slides: int = 5,
    sections: Optional[List[str]] = None,
    kpi_metrics: Optional[Dict[str, Any]] = None,
    qualitative_data: Optional[Dict[str, Any]] = None,
    botmaker_data: Optional[Dict[str, Any]] = None,
    yeastar_data: Optional[Dict[str, Any]] = None
) -> io.BytesIO:
    """
    Genera una presentación PowerPoint en memoria (BytesIO) completamente determinística.
    Soporta personalización de temas cromáticos, número de slides y secciones.
    """
    colors = THEMES.get(theme.lower(), THEMES["teal"])
    kpi_metrics = kpi_metrics or {}
    qualitative_data = qualitative_data or {}
    botmaker_data = botmaker_data or {}
    yeastar_data = yeastar_data or {}
    sections = sections or ["resumen", "eficiencia", "qualitative"]

    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Wide Aspect Ratio
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Slide 1 — Portada
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = colors["secondary"]
    bg1.line.color.rgb = colors["secondary"]

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"Reporte de Gestión & Calidad\n{client_name}"
    p1.font.bold = True
    p1.font.size = Pt(40)
    p1.font.color.rgb = colors["text_light"]
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"\nPeríodo: {period}  •  RC506 Quality & Operations System"
    p2.font.size = Pt(20)
    p2.font.color.rgb = colors["accent"]
    p2.alignment = PP_ALIGN.CENTER

    slides_created = 1

    # Helper para agregar encadenamiento de slide con título
    def add_standard_slide(title_text: str):
        slide = prs.slides.add_slide(blank_layout)
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = colors["primary"]
        header.line.color.rgb = colors["primary"]

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = colors["text_light"]
        return slide

    # Slide 2 — Resumen Ejecutivo (Si entra en num_slides y se solicita)
    if slides_created < num_slides and "resumen" in sections:
        s2 = add_standard_slide("📊 Resumen Ejecutivo de Métricas")
        kpi_cards = [
            ("Chats Gestionados", str(kpi_metrics.get("chats", 0)), "💬"),
            ("Leads Generados", str(kpi_metrics.get("leads", 0)), "🎯"),
            ("Ventas Cerradas", str(kpi_metrics.get("sales", 0)), "💰"),
            ("CSAT Promedio", f"{kpi_metrics.get('csat', 4.5)} / 5.0", "⭐")
        ]

        left_positions = [Inches(1.0), Inches(4.0), Inches(7.0), Inches(10.0)]
        for idx, (label, val, icon) in enumerate(kpi_cards):
            card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[idx], Inches(2.5), Inches(2.5), Inches(3.5))
            card.fill.solid()
            card.fill.fore_color.rgb = colors["bg"]
            card.line.color.rgb = colors["accent"]

            tb = s2.shapes.add_textbox(left_positions[idx], Inches(2.8), Inches(2.5), Inches(3.0))
            tf = tb.text_frame
            tf.word_wrap = True

            p_icon = tf.paragraphs[0]
            p_icon.text = icon
            p_icon.font.size = Pt(36)
            p_icon.alignment = PP_ALIGN.CENTER

            p_val = tf.add_paragraph()
            p_val.text = f"\n{val}"
            p_val.font.bold = True
            p_val.font.size = Pt(28)
            p_val.font.color.rgb = colors["primary"]
            p_val.alignment = PP_ALIGN.CENTER

            p_lbl = tf.add_paragraph()
            p_lbl.text = f"\n{label}"
            p_lbl.font.size = Pt(14)
            p_lbl.font.color.rgb = colors["text_dark"]
            p_lbl.alignment = PP_ALIGN.CENTER

        slides_created += 1

    # Slide 3 — Métricas Botmaker (Si entra en num_slides)
    if slides_created < num_slides and ("botmaker" in sections or "eficiencia" in sections):
        s3 = add_standard_slide("🤖 Operación Omnicanal Botmaker")
        bm = botmaker_data or {}
        tot_conv = bm.get("total_conversations") or kpi_metrics.get("chats", 0)
        tot_ag = bm.get("conversations_with_agent") or kpi_metrics.get("support", 0)
        tot_bot = bm.get("conversations_bot_only") or max(0, tot_conv - tot_ag)

        bm_items = [
            ("Conversaciones Totales", str(tot_conv)),
            ("Atendidas por Agente", str(tot_ag)),
            ("Resueltas por Bot", str(tot_bot)),
            ("Mensajes Usuario", str(bm.get("total_messages_user", 0))),
            ("Mensajes Bot", str(bm.get("total_messages_bot", 0))),
            ("Plantillas Enviadas", str(bm.get("total_templates_sent", 0)))
        ]

        for i, (l, v) in enumerate(bm_items):
            col = i % 3
            row = i // 3
            left = Inches(1.0 + col * 3.8)
            top = Inches(2.2 + row * 2.3)

            box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.5), Inches(1.8))
            box.fill.solid()
            box.fill.fore_color.rgb = colors["bg"]
            box.line.color.rgb = colors["primary"]

            tb = s3.shapes.add_textbox(left, top + Inches(0.2), Inches(3.5), Inches(1.4))
            tf = tb.text_frame
            p_v = tf.paragraphs[0]
            p_v.text = v
            p_v.font.bold = True
            p_v.font.size = Pt(26)
            p_v.font.color.rgb = colors["secondary"]
            p_v.alignment = PP_ALIGN.CENTER

            p_l = tf.add_paragraph()
            p_l.text = l
            p_l.font.size = Pt(14)
            p_l.font.color.rgb = colors["text_dark"]
            p_l.alignment = PP_ALIGN.CENTER

        slides_created += 1

    # Slide 4 — Análisis Cualitativo (Si entra en num_slides y se solicita)
    if slides_created < num_slides and "qualitative" in sections:
        s4 = add_standard_slide("📝 Análisis Cualitativo Estratégico")
        q_cards = [
            ("🔴 Puntos Críticos", qualitative_data.get("critical_points") or "Sin observaciones registradas.", RGBColor(239, 68, 68)),
            ("🟡 Advertencias", qualitative_data.get("warnings") or "Sin observaciones registradas.", RGBColor(245, 158, 11)),
            ("🟢 Logros / Positivo", qualitative_data.get("achievements") or "Sin observaciones registradas.", RGBColor(16, 185, 129)),
            ("🔵 Información General", qualitative_data.get("general_info") or "Sin observaciones registradas.", RGBColor(59, 130, 246))
        ]

        for i, (title, content, border_color) in enumerate(q_cards):
            col = i % 2
            row = i // 2
            left = Inches(1.0 + col * 5.8)
            top = Inches(1.8 + row * 2.6)

            card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(2.3))
            card.fill.solid()
            card.fill.fore_color.rgb = colors["bg"]
            card.line.color.rgb = border_color
            card.line.width = Pt(2)

            tb = s4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), Inches(4.9), Inches(2.1))
            tf = tb.text_frame
            tf.word_wrap = True

            p_t = tf.paragraphs[0]
            p_t.text = title
            p_t.font.bold = True
            p_t.font.size = Pt(16)
            p_t.font.color.rgb = colors["text_dark"]

            p_c = tf.add_paragraph()
            p_c.text = content[:250] + ("..." if len(content) > 250 else "")
            p_c.font.size = Pt(12)
            p_c.font.color.rgb = colors["text_dark"]

        slides_created += 1

    # Slide 5/6 — Cierre y Contacto (Siempre generado para completar la presentación)
    s_close = prs.slides.add_slide(blank_layout)
    bg_c = s_close.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg_c.fill.solid()
    bg_c.fill.fore_color.rgb = colors["secondary"]

    tb_c = s_close.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(3.0))
    tf_c = tb_c.text_frame
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "RC506 — Control de Calidad & Métricas"
    p_c1.font.bold = True
    p_c1.font.size = Pt(36)
    p_c1.font.color.rgb = colors["text_light"]
    p_c1.alignment = PP_ALIGN.CENTER

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "\nGracias por su confianza.\nInforme generado automáticamente por la Plataforma KAIROS."
    p_c2.font.size = Pt(18)
    p_c2.font.color.rgb = colors["accent"]
    p_c2.alignment = PP_ALIGN.CENTER

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
